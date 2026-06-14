import logging
import uuid

from django.contrib import messages
from django.contrib.auth.decorators import login_required
from django.db import transaction
from django.http import JsonResponse
from django.shortcuts import get_object_or_404, redirect, render
from django.utils import timezone

from users.decorators import role_required
from users.models import User
from .forms import StudyMaterialUploadForm, VideoUploadForm
from .models import Video

logger = logging.getLogger(__name__)


def _try_signed_url(key):
    if not key:
        return None
    try:
        from utils.r2_storage import get_signed_url
        return get_signed_url(key)
    except Exception:
        return None


def _course_is_locked(course, user):
    """Return True when an instructor course already has 1 video + 1 study material."""
    if user.role == User.Role.ADMIN:
        return False
    video_qs = Video.objects.filter(course=course)
    if not video_qs.exists():
        return False
    return bool(video_qs.first().english_transcript)


def _user_can_access_course(user, course):
    if user.role == User.Role.ADMIN:
        return True
    if user.role == User.Role.INSTRUCTOR:
        return course.instructor_id == user.id
    if user.role != User.Role.STUDENT:
        return False

    from courses.models import Enrollment
    return Enrollment.objects.filter(
        student=user, course=course, is_active=True
    ).exists()


def _user_can_manage_course(user, course):
    return user.role == User.Role.ADMIN or (
        user.role == User.Role.INSTRUCTOR and course.instructor_id == user.id
    )


def _repair_stale_video_status(video):
    """Backfill old status values and enforce processing timeout."""
    if video.processing_started_at is None and video.status == Video.Status.PROCESSING:
        fallback_started_at = video.created_at
        if fallback_started_at and video.has_study_material():
            video.processing_started_at = fallback_started_at
            video.save(update_fields=['processing_started_at'])
    return video.sync_runtime_status()


@role_required(User.Role.ADMIN, User.Role.INSTRUCTOR, message='Trainers and admins only.')
def upload_video_view(request):
    if request.method == 'POST':
        form = VideoUploadForm(request.POST, request.FILES, user=request.user)
        if form.is_valid():
            course = form.cleaned_data['course']
            if not _user_can_manage_course(request.user, course):
                messages.error(request, 'You do not have permission to manage that course.')
                return redirect('course-list')

            # Instructors: max 1 video per course
            if request.user.role == User.Role.INSTRUCTOR:
                if Video.objects.filter(course=course).exists():
                    messages.error(
                        request,
                        'Content already exists for this course. '
                        'Each course can have only one content item.'
                    )
                    return render(request, 'videos/upload.html', {'form': form})

            video = form.save(commit=False)
            video_file = request.FILES.get('video_file')

            if video_file:
                ext = (
                    video_file.name.rsplit('.', 1)[-1].lower()
                    if '.' in video_file.name else 'mp4'
                )
                key = f'videos/{course.id}/{uuid.uuid4()}.{ext}'
                try:
                    from utils.r2_storage import upload_file
                    upload_file(video_file, key)
                except Exception as exc:
                    logger.exception('R2 upload failed for course %s: %s', course.id, exc)
                    messages.error(
                        request,
                        'The video upload could not be completed right now. Please try again.'
                    )
                    return render(request, 'videos/upload.html', {'form': form})
                video.video_key = key

            video.status = Video.Status.UPLOADED
            video.save()

            messages.success(
                request,
                f'"{video.title}" saved. Now upload the study material to complete the content.'
            )
            return redirect('upload-material', video_id=video.id)
    else:
        form = VideoUploadForm(user=request.user)
    return render(request, 'videos/upload.html', {'form': form})


@role_required(User.Role.ADMIN, User.Role.INSTRUCTOR, message='Trainers and admins only.')
def upload_material_view(request, video_id):
    video = get_object_or_404(Video.objects.select_related('course'), id=video_id)
    if not _user_can_manage_course(request.user, video.course):
        messages.error(request, 'You do not have permission to manage that course.')
        return redirect('course-list')

    # Instructors: block re-upload once course is locked
    if request.user.role == User.Role.INSTRUCTOR:
        if _course_is_locked(video.course, request.user):
            messages.error(
                request,
                'This course already has a video and study material. '
                'Content cannot be modified once both are uploaded.'
            )
            return redirect('video-detail', video_id=video.id)

    if request.method == 'POST':
        form = StudyMaterialUploadForm(request.POST, request.FILES)
        if form.is_valid():
            material_file = request.FILES.get('material_file')
            english_text  = form.cleaned_data.get('english_content', '').strip()
            malayalam_text = form.cleaned_data.get('malayalam_content', '').strip()

            if material_file:
                try:
                    english_text = _extract_text(material_file)
                except Exception as exc:
                    logger.warning(
                        'Study material extraction failed for video %s: %s',
                        video.id,
                        exc,
                    )
                    messages.error(
                        request,
                        'The uploaded file could not be processed. '
                        'Please upload a valid PDF, DOCX, PPT, PPTX, or TXT file.'
                    )
                    return render(request, 'videos/upload_material.html', {
                        'form': form, 'video': video,
                    })

            video.english_transcript = english_text
            video.malayalam_transcript = malayalam_text
            video.status = Video.Status.PROCESSING
            video.processing_started_at = timezone.now()
            video.save(update_fields=[
                'english_transcript',
                'malayalam_transcript',
                'status',
                'processing_started_at',
            ])

            processing_queued = False
            try:
                from videos.tasks import process_study_material
                transaction.on_commit(lambda: process_study_material.delay(video.id))
                processing_queued = True
            except Exception as exc:
                logger.exception(
                    'Background study-material processing could not be queued for video %s: %s',
                    video.id,
                    exc,
                )
                messages.warning(
                    request,
                    'Study material was saved, but background processing is temporarily unavailable.'
                )

            if processing_queued:
                messages.success(
                    request,
                    'Study material saved. Quiz draft generation and content processing started in the background.'
                )

            return redirect('video-detail', video_id=video.id)
    else:
        form = StudyMaterialUploadForm(initial={
            'english_content':   video.english_transcript,
            'malayalam_content': video.malayalam_transcript,
        })

    return render(request, 'videos/upload_material.html', {'form': form, 'video': video})


# ---------------------------------------------------------------------------
# File text extraction helpers
# ---------------------------------------------------------------------------

def _extract_text(uploaded_file):
    name = uploaded_file.name.lower()
    data = uploaded_file.read()

    if name.endswith('.txt'):
        return data.decode('utf-8', errors='replace').strip()
    if name.endswith('.pdf'):
        return _extract_pdf(data)
    if name.endswith('.docx'):
        return _extract_docx(data)
    if name.endswith(('.pptx', '.ppt')):
        return _extract_pptx(data)

    raise ValueError(f'Unsupported file type: {uploaded_file.name}')


def _extract_pdf(data):
    import io as _io
    from pypdf import PdfReader
    reader = PdfReader(_io.BytesIO(data))
    pages = []
    for page in reader.pages:
        text = page.extract_text() or ''
        if text.strip():
            pages.append(text.strip())
    return '\n\n'.join(pages)


def _extract_docx(data):
    import io as _io
    from docx import Document
    doc = Document(_io.BytesIO(data))
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return '\n\n'.join(paras)


def _extract_pptx(data):
    import io as _io
    from pptx import Presentation
    prs = Presentation(_io.BytesIO(data))
    slides = []
    for slide in prs.slides:
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, 'text') and shape.text.strip()
        ]
        if texts:
            slides.append('\n'.join(texts))
    return '\n\n'.join(slides)


@login_required
def video_list_view(request, course_id):
    from courses.models import Course
    course = get_object_or_404(Course, id=course_id)
    if not _user_can_access_course(request.user, course):
        messages.error(
            request,
            'You do not have permission to view that course content.'
        )
        return redirect('course-list')
    videos = Video.objects.filter(course=course)
    return render(request, 'videos/video_list.html', {'course': course, 'videos': videos})


@login_required
def video_detail_view(request, video_id):
    video = _repair_stale_video_status(
        get_object_or_404(Video.objects.select_related('course'), id=video_id)
    )

    if not _user_can_access_course(request.user, video.course):
        messages.error(
            request,
            'You do not have permission to view that content.'
        )
        return redirect('course-list')

    course_locked = _course_is_locked(video.course, request.user)

    chatbot_enrolled = request.user.role == User.Role.STUDENT and _user_can_access_course(
        request.user, video.course
    )

    pending_draft_count = video.quiz_drafts.filter(status='pending').count()
    approved_draft_count = video.quiz_drafts.filter(status='approved').count()
    chunks = video.chunks.all()
    return render(request, 'videos/video_detail.html', {
        'video':             video,
        'chunks':            chunks,
        'course_locked':     course_locked,
        'pending_draft_count': pending_draft_count,
        'approved_draft_count': approved_draft_count,
        'processing_timeout_minutes': int(Video.PROCESSING_TIMEOUT.total_seconds() // 60),
        'signed_url':        _try_signed_url(video.video_key),
        'english_pdf_url':   _try_signed_url(video.english_pdf_key),
        'malayalam_pdf_url': _try_signed_url(video.malayalam_pdf_key),
        'chatbot_course_id': video.course_id if chatbot_enrolled else 0,
        'chatbot_courses':   (
            [{'id': video.course_id, 'title': video.course.title}]
            if chatbot_enrolled else []
        ),
        'chatbot_enrolled':  chatbot_enrolled,
    })


@role_required(User.Role.ADMIN, User.Role.INSTRUCTOR, message='Trainers and admins only.')
def generate_quiz_view(request, video_id):
    video = get_object_or_404(Video.objects.select_related('course'), id=video_id)
    if not _user_can_manage_course(request.user, video.course):
        messages.error(request, 'You do not have permission to manage that content.')
        return redirect('course-list')

    if request.method != 'POST':
        return redirect('video-detail', video_id=video.id)

    if not (video.english_transcript or '').strip():
        messages.error(request, 'Quiz generation requires study material first.')
        return redirect('video-detail', video_id=video.id)

    try:
        from videos.tasks import queue_quiz_generation
        queued = queue_quiz_generation(video.id)
        if not queued:
            pending_draft_count = video.quiz_drafts.filter(status='pending').count()
            approved_draft_count = video.quiz_drafts.filter(status='approved').count()
            if pending_draft_count or approved_draft_count:
                messages.info(
                    request,
                    f'Quiz draft already exists for this content. Open Quiz Drafts to review {pending_draft_count} pending and {approved_draft_count} approved question(s).'
                )
                return redirect('quiz-draft-list')
            has_published_quiz = video.quizzes.filter(is_published=True).exists()
            has_unpublished_quiz = video.quizzes.filter(is_published=False).exists()
            if has_published_quiz:
                messages.info(
                    request,
                    'Quiz generation was skipped because this content already has a published quiz.'
                )
            elif has_unpublished_quiz:
                messages.info(
                    request,
                    'Quiz generation was skipped because this content already has a saved quiz.'
                )
            else:
                messages.info(
                    request,
                    'Quiz generation was skipped because this content already has a quiz.'
                )
            return redirect('video-detail', video_id=video.id)
        messages.success(
            request,
            'Quiz generation started. The draft will appear in the review queue shortly.'
        )
    except Exception as exc:
        logger.exception('Quiz generation queue failed for video %s: %s', video.id, exc)
        messages.warning(
            request,
            'Quiz generation is temporarily unavailable. Please try again shortly.'
        )

    return redirect('video-detail', video_id=video.id)


@login_required
def video_status_api(request, video_id):
    video = _repair_stale_video_status(
        get_object_or_404(Video.objects.select_related('course'), id=video_id)
    )
    if not _user_can_access_course(request.user, video.course):
        return JsonResponse({'error': 'Forbidden'}, status=403)
    return JsonResponse({
        'id':                video.id,
        'status':            video.status,
        'status_display':    video.get_status_display(),
        'chunk_count':       video.chunks.count(),
        'has_english_pdf':   bool(video.english_pdf_key),
        'has_malayalam_pdf': bool(video.malayalam_pdf_key),
        'processing_timed_out': video.status == Video.Status.FAILED,
    })
